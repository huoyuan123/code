import sys
from pathlib import Path
from pptx import Presentation

def extract_text_from_pptx(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    out_lines = []
    for si, slide in enumerate(prs.slides, start=1):
        out_lines.append(f"## Slide {si}")
        for shape in slide.shapes:
            # Text frames
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Join all runs to preserve emphasis content as plain text
                    txt = "".join(run.text for run in para.runs)
                    txt = txt.strip()
                    if txt:
                        bullet = "- " if para.level >= 0 else ""
                        out_lines.append(f"{bullet}{txt}")
            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for r in table.rows:
                    row_txt = []
                    for c in r.cells:
                        row_txt.append(c.text_frame.text.strip())
                    if any(row_txt):
                        out_lines.append(" | ".join(row_txt))
        out_lines.append("")
    return "\n".join(out_lines)


def main():
    if len(sys.argv) < 3:
        print("Usage: python extract_pptx_text.py <input.pptx> <output.txt>")
        sys.exit(1)
    in_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    if not in_path.exists():
        print(f"Input not found: {in_path}")
        sys.exit(2)
    text = extract_text_from_pptx(in_path)
    out_path.write_text(text, encoding="utf-8")
    print(f"Wrote text to {out_path}")

if __name__ == "__main__":
    main()
